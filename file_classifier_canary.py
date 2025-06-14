'''
pyinstaller -D -n "file_classifier" -i "icon.ico" --hidden-import "sklearn.utils._typedefs" --hidden-import "sklear
n.utils._cython_blas" --hidden-import "sklearn.neighbors.typedefs" --hidden-import "sklearn.neighbors.quad_tree" --hidden-import "sklearn.tree._util
s" --hidden-import "winrt.windows.ui.notifications" --hidden-import "winrt.windows.data.xml.dom" --hidden-import "pystray._win32" --hidden-import "p
ytesseract" --hidden-import "cv2" --collect-all "sklearn" --collect-all "joblib" --collect-all "threadpoolctl" --collect-all "scipy" --collect-all "numpy" --collect-all "PIL" file_classifier_canary.py -w
'''
import os
import sys
import time
import re
import shutil
import threading
from collections import defaultdict
import json
import logging
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from PIL import Image
from pystray import Icon as icon, Menu as menu, MenuItem as item
import joblib
import docx
import pptx
import fitz  # PyMuPDF
import pytesseract
import cv2
import tempfile
import subprocess
import winrt.windows.ui.notifications as notifications
import winrt.windows.data.xml.dom as dom
import tkinter as tk
from tkinter import messagebox

# ==================== 配置参数 ====================
CONFIG_FILE = "config.json"
try:
    with open(CONFIG_FILE, "r", encoding="utf-8") as f:
        config = json.load(f)
except FileNotFoundError:
    print("⚠️ 未找到配置文件，请先运行配置工具")
    sys.exit(1)

WATCH_FOLDER = config["WATCH_FOLDER"]
OUTPUT_BASE_FOLDER = config["OUTPUT_BASE_FOLDER"]
SUPPORTED_EXTS = config["SUPPORTED_EXTS"]
DELAY_SECONDS = config.get("DELAY_SECONDS", 3)

pytesseract.pytesseract.tesseract_cmd = config["TESSERACT_PATH"]

# 学科关键词映射表
SUBJECT_KEYWORDS = {
    '语文': ['语文', '文言文', '古诗', '作文', '阅读理解', '现代文', '议论文', '小说鉴赏', '散文赏析', '写作指导'],
    '数学': ['数学', '函数', '几何', '方程', '代数', '集合', '不等式', '三角函数', '概率统计', '解析几何', '立体几何'],
    '英语': ['英语', '单词', '语法', 'reading', '听力', '作文', '完形填空', '阅读理解', '语法填空', '短文改错'],
    '物理': ['ck', '物理', '力', '能量', '电学', '磁场', '电磁感应', '动量守恒', '牛顿定律', '曲线运动', '机械振动'],
    '化学': ['化学', '反应', '元素', '物质', '方程式', '离子反应', '氧化还原', '有机化学', '无机化学', '化学平衡', '化合物'],
    '生物': ['生物', '细胞', 'DNA', '生态', '遗传', '基因', '光合作用', '呼吸作用', '种群', '生态系统'],
    '历史': ['历史', '朝代', '战争', '文明', '近代史', '古代史', '世界史', '中国史', '改革', '政治制度'],
    '政治': ['政治', '法律', '公民', '制度', '国家', '人民代表大会', '中国特色社会主义', '哲学生活', '文化与生活'],
    '地理': ['地理', '气候', '地形', '区域', '地图', '地球运动', '大气环流', '水循环', '农业区位', '工业区位']
}

# ==================== 工具函数 ====================
def clean_text(text):
    return re.sub(r'\s+', ' ', text).strip()

def ensure_folder_exists(folder):
    if not os.path.exists(folder):
        os.makedirs(folder)

# ==================== 提取内容函数 ====================
def extract_docx(file_path):
    try:
        doc = docx.Document(file_path)
        return '\n'.join(para.text for para in doc.paragraphs)
    except Exception as e:
        log(f"DOCX提取失败: {file_path} - {e}")
        return None

def extract_pptx(file_path):
    try:
        pres = pptx.Presentation(file_path)
        return '\n'.join(shape.text for slide in pres.slides
                         for shape in slide.shapes if hasattr(shape, "text"))
    except Exception as e:
        log(f"PPTX提取失败: {file_path} - {e}")
        return None

def extract_pdf(file_path):
    try:
        doc = fitz.open(file_path)
        text = '\n'.join(page.get_text() for page in doc)
        if not text.strip():
            raise ValueError("PDF无文本内容，尝试OCR")
        return text
    except Exception as e:
        log(f"PDF原生提取失败，尝试OCR：{file_path} - {e}")
        return ocr_pdf(file_path)

def ocr_image(image_path):
    try:
        img = cv2.imread(image_path)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        _, binary = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)
        text = pytesseract.image_to_string(binary, lang='chi_sim')
        return text
    except Exception as e:
        log(f"OCR识别失败: {image_path} - {e}")
        return ""

def ocr_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        temp_dir = tempfile.mkdtemp()
        texts = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=200)
            img_path = os.path.join(temp_dir, f"page_{page_num}.png")
            pix.save(img_path)
            texts.append(ocr_image(img_path))

        shutil.rmtree(temp_dir)
        return '\n'.join(texts)
    except Exception as e:
        log(f"OCR转换失败: {pdf_path} - {e}")
        return ""

def extract_content(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    elif ext == ".pdf":
        return extract_pdf(file_path)
    else:
        return ""

# ==================== 加载模型 ====================
try:
    clf = joblib.load("subject_classifier.pkl")
    vectorizer = joblib.load("tfidf_vectorizer.pkl")
except Exception as e:
    print("⚠️ 模型文件未找到，请先运行训练脚本！", e)
    sys.exit(1)

# ==================== 文件名辅助判断 ====================
def guess_by_filename(filename):
    filename_base = os.path.splitext(os.path.basename(filename))[0]
    filename_base = clean_text(filename_base)
    scores = defaultdict(int)

    for subject, keywords in SUBJECT_KEYWORDS.items():
        for keyword in keywords:
            if re.search(keyword, filename_base, re.IGNORECASE):
                scores[subject] += 1

    if not scores:
        return None

    max_score = max(scores.values())
    candidates = [k for k, v in scores.items() if v == max_score]
    return candidates[0] if len(candidates) == 1 else None

# ==================== 主分类逻辑（先文件名后内容）====================
def classify_file(file_path):
    subject = guess_by_filename(file_path)
    if subject:
        log(f"[文件名识别] 文件 {os.path.basename(file_path)} 分类为：{subject}")
        return subject

    content = extract_content(file_path)
    if content and len(content.strip()) >= 30:
        X = vectorizer.transform([content])
        ai_subject = clf.predict(X)[0]
        log(f"[AI模型识别] 文件 {os.path.basename(file_path)} 分类为：{ai_subject}")
        return ai_subject

    return None

# ==================== 日志模块 ====================
logger = logging.getLogger()
logger.setLevel(logging.INFO)
handler = logging.FileHandler("file_classifier.log", mode='a', encoding='utf-8')
formatter = logging.Formatter('%(asctime)s - %(message)s')
handler.setFormatter(formatter)
logger.addHandler(handler)

def log(msg):
    print(f"[{time.strftime('%D  %H:%M:%S')}] {msg}")
    logger.info(msg)

# ==================== 文件处理缓存 ====================
processed_files = set()
processed_lock = threading.Lock()


# ===================== 通知 =========================
def save_result(result):
    with open('temp_msg.txt', 'w', encoding='utf-8') as file:
        file.write(result)

def send_toast(subject, file_path):
    if file_path == "STARTUP":
        title = "AI课件分类器已启动"
        message = f"开始监视文件夹: {WATCH_FOLDER}\n文件将分类到: {OUTPUT_BASE_FOLDER}"
    else:
        title = f"分类：{subject}"
        message = f"文件 {os.path.basename(file_path)} 分类为：{subject}\n是否正确？（请在5s内判断）\n如果正确请不要理睬此消息，或关闭\n如果错误，请点击此消息，文件会自动移回"

    toast_xml = notifications.ToastNotificationManager.get_template_content(
        notifications.ToastTemplateType.TOAST_TEXT02)

    text_elements = toast_xml.get_elements_by_tag_name("text")
    text_elements[0].append_child(toast_xml.create_text_node(title))
    text_elements[1].append_child(toast_xml.create_text_node(message))

    notifier = notifications.ToastNotificationManager.create_toast_notifier("AI课件分类器")
    toast = notifications.ToastNotification(toast_xml)

    if file_path != "STARTUP":  # 只有分类通知才需要处理反馈
        toast.add_activated(lambda _, __: save_result("no"))
        toast.add_dismissed(lambda _, __: save_result("yes"))

    notifier.show(toast)


# ==================== 文件监控类 ====================
class FileHandler(FileSystemEventHandler):
    def __init__(self):
        self.running = True

    def on_created(self, event):
        if not self.running or event.is_directory:
            return

        file_path = event.src_path
        _, ext = os.path.splitext(file_path)
        if ext.lower() not in SUPPORTED_EXTS:
            log(f"跳过不支持的文件类型：{file_path}")
            return

        with processed_lock:
            if file_path in processed_files:
                log(f"跳过已处理文件：{file_path}")
                return
            processed_files.add(file_path)

        threading.Thread(target=self.delayed_classification, args=(file_path,), daemon=True).start()

    def delayed_classification(self, file_path):
        log(f"等待 {DELAY_SECONDS} 秒后尝试识别文件：{file_path}")
        time.sleep(DELAY_SECONDS)

        if not self.running or not os.path.exists(file_path):
            log(f"文件已消失或监视已暂停: {file_path}")
            with processed_lock:
                if file_path in processed_files:
                    processed_files.remove(file_path)
            return

        try:
            with open(file_path, 'rb'):
                pass
        except Exception as e:
            log(f"文件无法访问（仍在复制中？）：{file_path} - {e}")
            return

        subject = classify_file(file_path)
        filename = os.path.basename(file_path)

        if subject:
            dest_folder = os.path.join(OUTPUT_BASE_FOLDER, subject)
            ensure_folder_exists(dest_folder)
            dest_path = os.path.join(dest_folder, filename)
            try:
                if not self.handle_user_feedback(subject, file_path):
                    log(f"用户标记为误判，文件保留在原始位置")
                else:
                    shutil.move(file_path, dest_path)
                    log(f"文件 {filename} 分类为：{subject}，已归类到 {dest_folder}/")
                    try:
                        os.remove('temp_msg.txt')
                    except FileNotFoundError:
                        pass
            except Exception as e:
                log(f"文件处理失败：{filename} - {e}")
        else:
            log(f"无法识别文件类型或内容太少：{filename}")
            with processed_lock:
                if file_path in processed_files:
                    processed_files.remove(file_path)

    def handle_user_feedback(self, subject, file_path):
        log(f"等待用户反馈: {os.path.basename(file_path)}")
        send_toast(subject, file_path)

        time.sleep(5)
        try:
            with open('temp_msg.txt', 'r', encoding='utf-8') as f:
                user_feedback = f.read().strip()
        except FileNotFoundError:
            user_feedback = "yes"

        if user_feedback == "no":
            log("用户反馈: 分类错误")
            move_back(file_path)
            return False
        else:
            log("用户反馈: 分类正确或超时")
            return True

# ==================== 用户反馈处理 ====================
def move_back(file_path):
    filename = os.path.basename(file_path)
    log(f"文件 {filename} 被用户标记为误判，保留在原始位置")

    with processed_lock:
        if file_path in processed_files:
            processed_files.remove(file_path)
            log(f"已从处理缓存中移除: {filename}")

    try:
        os.remove('temp_msg.txt')
    except FileNotFoundError:
        pass

# ==================== 启动后台监控 ====================
watcher = None

def start_file_watcher():
    global watcher
    log('监视文件夹:' + WATCH_FOLDER)
    log('移动文件夹:' + OUTPUT_BASE_FOLDER)

    if watcher is None:
        watcher = Observer()
        event_handler = FileHandler()
        watcher.schedule(event_handler, path=WATCH_FOLDER, recursive=False)
        watcher.start()
        log("文件监视已启动")

        # 发送启动通知
        try:
            send_toast("", "STARTUP")
        except Exception as e:
            log(f"发送启动通知失败: {e}")
    else:
        log("文件监视已在运行")

def stop_file_watcher():
    global watcher
    if watcher is not None:
        watcher.stop()
        watcher.join()
        watcher = None
        log("文件监视已暂停")
    else:
        log("文件监视未启动")

# ==================== 托盘图标逻辑 ====================
def create_tray_icon():
    return Image.open("icon.ico")

def view_log(icon, item):
    if os.path.exists("file_classifier.log"):
        try:
            os.startfile("file_classifier.log")
        except:
            opener = "open" if sys.platform == "darwin" else "xdg-open"
            subprocess.call([opener, "file_classifier.log"])
    else:
        messagebox.showerror("错误", "日志文件不存在")

def clear_log(icon, item):
    try:
        open("file_classifier.log", "w").close()
        messagebox.showinfo("提示", "日志已清空")
    except Exception as e:
        messagebox.showerror("错误", f"清空日志失败: {e}")

def open_config(icon, item):
    config_editor_path = "config_editor.py"
    if os.path.exists("config_editor.exe"):
        config_editor_path = "config_editor.exe"

    if not os.path.exists(config_editor_path):
        messagebox.showerror("错误", "找不到配置编辑器")
        return

    try:
        subprocess.Popen(
            [sys.executable, config_editor_path] if config_editor_path.endswith('.py') else [config_editor_path]
        )
    except Exception as e:
        messagebox.showerror("错误", f"启动配置界面失败：{e}")

def toggle_watcher(icon, item):
    if watcher is None:
        start_file_watcher()
        item.text = "暂停监视"
    else:
        stop_file_watcher()
        item.text = "恢复监视"

def exit_app(icon, item):
    if watcher is not None:
        stop_file_watcher()
    icon.stop()
    os._exit(0)

# ==================== 主程序入口 ====================
if __name__ == "__main__":
    start_file_watcher()

    tray_icon = icon("AI课件分类器", create_tray_icon(), menu=menu(
        item('打开配置界面', open_config),
        item('查看日志', view_log),
        item('清空日志', clear_log),
        item('暂停/恢复监视', toggle_watcher, default=True),
        item('退出', exit_app)
    ))


    tray_icon.run()