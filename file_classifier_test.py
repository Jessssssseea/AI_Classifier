import os
import sys
import time
import re
import shutil
import threading
from collections import defaultdict
import json
import logging
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from PIL import Image, ImageDraw
from pystray import Icon as icon, Menu as menu, MenuItem as item
import joblib
import docx
import pptx
import fitz  # PyMuPDF
import pytesseract
import cv2
import tempfile
from tkinter import Tk, messagebox

# ==================== 配置参数 ====================
CONFIG_FILE = "config.json"
try:
    with open(CONFIG_FILE, "r", encoding="utf-8") as f:
        config = json.load(f)
except FileNotFoundError:
    print("⚠️ 未找到配置文件，请先运行配置工具")
    exit(1)

WATCH_FOLDER = config["WATCH_FOLDER"]
OUTPUT_BASE_FOLDER = config["OUTPUT_BASE_FOLDER"]
SUPPORTED_EXTS = config["SUPPORTED_EXTS"]
DELAY_SECONDS = config["DELAY_SECONDS"]

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# 学科关键词映射表（保留不变）
SUBJECT_KEYWORDS = {
    '语文': ['语文', '文言文', '古诗', '作文', '阅读理解', '现代文', '议论文', '小说鉴赏', '散文赏析', '写作指导'],
    '数学': ['数学', '函数', '几何', '方程', '代数', '集合', '不等式', '三角函数', '概率统计', '解析几何', '立体几何'],
    '英语': ['英语', '单词', '语法', 'reading', '听力', '作文', '完形填空', '阅读理解', '语法填空', '短文改错'],
    '物理': ['ck', '物理', '力', '能量', '电学', '磁场', '电磁感应', '动量守恒', '牛顿定律', '曲线运动', '机械振动'],
    '化学': ['化学', '反应', '元素', '物质', '方程式', '离子反应', '氧化还原', '有机化学', '无机化学', '化学平衡', '化合物'],
    '生物': ['生物', '细胞', 'DNA', '生态', '遗传', '基因', '光合作用', '呼吸作用', '种群', '生态系统'],
    '历史': ['历史', '朝代', '战争', '文明', '近代史', '古代史', '世界史', '中国史', '改革', '政治制度'],
    '政治': ['政治', '法律', '公民', '制度', '国家', '人民代表大会', '中国特色社会主义', '哲学生活', '文化与生活'],
    '地理': ['地理', '气候', '地形', '区域', '地图', '地球运动', '大气环流', '水循环', '农业区位', '工业区位']
}

# ==================== 工具函数 ====================
def clean_text(text):
    return re.sub(r'\s+', ' ', text).strip()

def ensure_folder_exists(folder):
    if not os.path.exists(folder):
        os.makedirs(folder)

# ==================== 提取内容函数 ====================
def extract_docx(file_path):
    try:
        doc = docx.Document(file_path)
        return '\n'.join(para.text for para in doc.paragraphs)
    except Exception as e:
        log(f"DOCX提取失败: {file_path} - {e}")
        return None

def extract_pptx(file_path):
    try:
        pres = pptx.Presentation(file_path)
        return '\n'.join(shape.text for slide in pres.slides
                        for shape in slide.shapes if hasattr(shape, "text"))
    except Exception as e:
        log(f"PPTX提取失败: {file_path} - {e}")
        return None

def extract_pdf(file_path):
    try:
        doc = fitz.open(file_path)
        text = '\n'.join(page.get_text() for page in doc)
        if not text.strip():
            raise ValueError("PDF无文本内容，尝试OCR")
        return text
    except Exception as e:
        log(f"PDF原生提取失败，尝试OCR：{file_path} - {e}")
        return ocr_pdf(file_path)

def ocr_image(image_path):
    try:
        img = cv2.imread(image_path)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        _, binary = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)
        text = pytesseract.image_to_string(binary, lang='chi_sim')
        return text
    except Exception as e:
        log(f"OCR识别失败: {image_path} - {e}")
        return ""

def ocr_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        temp_dir = tempfile.mkdtemp()
        texts = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=200)
            img_path = os.path.join(temp_dir, f"page_{page_num}.png")
            pix.save(img_path)
            texts.append(ocr_image(img_path))

        shutil.rmtree(temp_dir)
        return '\n'.join(texts)
    except Exception as e:
        log(f"OCR转换失败: {pdf_path} - {e}")
        return ""

def extract_content(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    elif ext == ".pdf":
        return extract_pdf(file_path)
    else:
        return ""

# ==================== 加载模型 ====================
try:
    clf = joblib.load("subject_classifier.pkl")
    vectorizer = joblib.load("tfidf_vectorizer.pkl")
except Exception as e:
    print("⚠️ 模型文件未找到，请先运行训练脚本！", e)
    exit(1)

# ==================== 文件名辅助判断（优先）====================
def guess_by_filename(filename):
    filename_base = os.path.splitext(os.path.basename(filename))[0]
    filename_base = clean_text(filename_base)
    scores = defaultdict(int)

    for subject, keywords in SUBJECT_KEYWORDS.items():
        for keyword in keywords:
            if re.search(keyword, filename_base, re.IGNORECASE):
                scores[subject] += 1

    if not scores:
        return None

    max_score = max(scores.values())
    candidates = [k for k, v in scores.items() if v == max_score]
    return candidates[0] if len(candidates) == 1 else None

# ==================== 主分类逻辑（先文件名后内容）====================
def classify_file(file_path):
    subject = guess_by_filename(file_path)
    if subject:
        log(f"[文件名识别] 文件 {os.path.basename(file_path)} 分类为：{subject}")
        return subject

    content = extract_content(file_path)
    if content and len(content.strip()) >= 30:
        X = vectorizer.transform([content])
        ai_subject = clf.predict(X)[0]
        log(f"[AI模型识别] 文件 {os.path.basename(file_path)} 分类为：{ai_subject}")
        return ai_subject

    return None

# ==================== 日志模块 ====================
logger = logging.getLogger()
logger.setLevel(logging.INFO)
handler = logging.FileHandler("file_classifier.log", mode='a', encoding='utf-8')
formatter = logging.Formatter('%(asctime)s - %(message)s')
handler.setFormatter(formatter)
logger.addHandler(handler)

def log(msg):
    print(f"[{time.strftime('%H:%M:%S')}] {msg}")
    logger.info(msg)

# ==================== 文件处理缓存 ====================
processed_files = set()
processed_lock = threading.Lock()

# ==================== 文件监控类 ====================
class FileHandler(FileSystemEventHandler):
    def on_created(self, event):
        if event.is_directory:
            return

        file_path = event.src_path
        _, ext = os.path.splitext(file_path)
        if ext.lower() not in SUPPORTED_EXTS:
            return

        with processed_lock:
            if file_path in processed_files:
                log(f"跳过已处理文件：{file_path}")
                return
            processed_files.add(file_path)

        threading.Thread(target=self.delayed_classification, args=(file_path,), daemon=True).start()

    def delayed_classification(self, file_path):
        log(f"等待 {DELAY_SECONDS} 秒后尝试识别文件：{file_path}")
        time.sleep(DELAY_SECONDS)

        if not os.path.exists(file_path):
            msg = f"文件不存在：{file_path}"
            log(msg)
            show_error_popup(file_path, msg)
            return

        try:
            with open(file_path, 'rb'):
                pass
        except Exception as e:
            msg = f"文件无法访问（仍在复制中？）：{file_path} - {e}"
            log(msg)
            show_error_popup(file_path, msg)
            return

        subject = classify_file(file_path)
        filename = os.path.basename(file_path)

        if subject:
            dest_folder = os.path.join(OUTPUT_BASE_FOLDER, subject)
            ensure_folder_exists(dest_folder)
            dest_path = os.path.join(dest_folder, filename)
            try:
                shutil.move(file_path, dest_path)
                msg = f"文件 {filename} 分类为：{subject}，已归类到 {dest_folder}/"
                log(msg)
                show_success_popup(dest_path, subject, dest_folder)
            except Exception as e:
                msg = f"文件移动失败：{filename} -> {dest_folder} - {e}"
                log(msg)
                show_error_popup(file_path, msg)
        else:
            msg = f"无法识别文件类型或内容太少：{filename}"
            log(msg)
            show_unknown_popup(file_path)

# ==================== GUI 弹窗功能 ====================
def show_success_popup(file_path, subject, dest_folder):
    root = Tk()
    root.withdraw()  # 隐藏主窗口
    response = messagebox.askokcancel(
        "文件分类成功",
        f"文件 {os.path.basename(file_path)} 已自动归类为：{subject}\n目标路径：{dest_folder}\n\n如果分类错误，请点击【取消】将文件移回原目录"
    )
    if not response:
        move_back(file_path)
    root.destroy()

def show_error_popup(file_path, message):
    root = Tk()
    root.withdraw()
    messagebox.showerror("错误", message)
    root.destroy()

def show_unknown_popup(file_path):
    root = Tk()
    root.withdraw()
    messagebox.showwarning("无法识别", f"无法识别文件 {os.path.basename(file_path)} 的学科\n请手动处理或重新命名后再拖入")
    root.destroy()

def move_back(file_path):
    original_path = file_path
    filename = os.path.basename(file_path)
    new_path = os.path.join(config["WATCH_FOLDER"], filename)

    if os.path.exists(new_path):
        try:
            os.remove(new_path)
        except Exception as e:
            log(f"删除旧文件失败：{new_path} - {e}")

    try:
        shutil.move(original_path, new_path)
        log(f"文件 {filename} 被标记为误判，已移回监控目录")
        messagebox.showinfo("操作完成", f"文件 {filename} 已移回监控目录")
    except Exception as e:
        log(f"文件移回失败：{filename} -> {config['WATCH_FOLDER']} - {e}")
        messagebox.showerror("错误", f"文件移回失败：{e}")

# ==================== 启动后台监控 ====================
def start_file_watcher():
    observer = Observer()
    observer.schedule(FileHandler(), path=config["WATCH_FOLDER"], recursive=False)
    observer.start()

    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()

# ==================== 托盘图标逻辑 ====================
def create_tray_icon():
    image = Image.new('RGB', (64, 64), color=(255, 255, 255))
    dc = ImageDraw.Draw(image)
    dc.rectangle((10, 10, 54, 54), fill=(0, 128, 255))
    return image

def view_log(icon, item):
    if os.path.exists("file_classifier.log"):
        try:
            os.startfile("file_classifier.log")
        except:
            opener = "open" if sys.platform == "darwin" else "xdg-open"
            subprocess.call([opener, "file_classifier.log"])

def clear_log(icon, item):
    try:
        open("file_classifier.log", "w").close()
        messagebox.showinfo("提示", "日志已清空")
    except Exception as e:
        messagebox.showerror("错误", f"清空日志失败: {e}")

def open_config(icon, item):
    config_editor_path = "config_editor.py"
    if os.path.exists("config_editor.exe"):
        config_editor_path = "config_editor.exe"

    if not os.path.exists(config_editor_path):
        messagebox.showerror("错误", "找不到配置编辑器")
        return

    try:
        subprocess.Popen([sys.executable, config_editor_path] if config_editor_path.endswith('.py') else [config_editor_path])
    except Exception as e:
        messagebox.showerror("错误", f"启动配置界面失败：{e}")

def exit_app(icon, item):
    icon.stop()
    os._exit(0)

# ==================== 主程序入口 ====================
if __name__ == "__main__":
    watcher_thread = threading.Thread(target=start_file_watcher, daemon=True)
    watcher_thread.start()

    tray_icon = icon("AI文件分类器", create_tray_icon(), menu=menu(
        item('打开配置界面', open_config),
        item('查看日志', view_log),
        item('清空日志', clear_log),
        item('退出', exit_app)
    ))

    print("应用已启动，常驻系统托盘中...")
    tray_icon.run()