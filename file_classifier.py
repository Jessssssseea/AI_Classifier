import os
import re
import time
import shutil
import threading
import win32api  # Windows API 获取U盘信息
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from plyer import notification
import joblib
import docx
import pptx
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import cv2
import tempfile
from collections import defaultdict

# ==================== 配置参数 ====================
WATCH_FOLDER = os.path.expanduser("D:\定期练手\自动归纳\monitor")  # 默认监控路径
OUTPUT_BASE_FOLDER = os.path.expanduser("D:\定期练手\自动归纳\monitor")  # 归类根目录
SUPPORTED_EXTS = ['.docx', '.pdf', '.pptx', '.wps', '.mp4', '.wbd']  # 支持的文件类型
DELAY_SECONDS = 3  # 延迟识别时间（秒）

# 设置 Tesseract 路径（如未加入系统环境变量）
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# U盘卷标映射表
USB_LABEL_MAP = {
    '黄': '语文',
    '杨': '数学',
    '黄漫霞备份': '英语'
}

# 科目关键词映射
SUBJECT_KEYWORDS = {
    '语文': ['语文', '文言文', '古诗', '作文', '阅读理解', '现代文', '议论文', '小说鉴赏', '散文赏析', '写作指导'],
    '数学': ['数学', '函数', '几何', '方程', '代数', '集合', '不等式', '三角函数', '概率统计', '解析几何', '立体几何'],
    '英语': ['英语', '单词', '语法', 'reading', '听力', '作文', '完形填空', '阅读理解', '语法填空', '短文改错'],
    '物理': ['ck', '物理', '力', '能量', '电学', '磁场', '电磁感应', '动量守恒', '牛顿定律', '曲线运动', '机械振动'],
    '化学': ['化学', '反应', '元素', '物质', '方程式', '离子反应', '氧化还原', '有机化学', '无机化学', '化学平衡', '化合物'],
    '生物': ['生物', '细胞', 'DNA', '生态', '遗传', '基因', '光合作用', '呼吸作用', '种群', '生态系统'],
    '历史': ['历史', '朝代', '战争', '文明', '近代史', '古代史', '世界史', '中国史', '改革', '政治制度'],
    '政治': ['政治', '法律', '公民', '制度', '国家', '人民代表大会', '中国特色社会主义', '哲学生活', '文化与生活'],
    '地理': ['地理', '气候', '地形', '区域', '地图', '地球运动', '大气环流', '水循环', '农业区位', '工业区位']
}


# ==================== 工具函数 ====================
def clean_text(text):
    return re.sub(r'\s+', ' ', text).strip()


def ensure_folder_exists(folder):
    if not os.path.exists(folder):
        os.makedirs(folder)


def get_usb_label(file_path):
    """获取文件所在磁盘卷标"""
    try:
        drive = os.path.splitdrive(file_path)[0][0]
        return win32api.GetVolumeInformation(f"{drive}:\\")[0]
    except:
        return None


# ==================== 提取内容函数 ====================
def extract_docx(file_path):
    try:
        doc = docx.Document(file_path)
        return '\n'.join(para.text for para in doc.paragraphs)
    except Exception as e:
        print(f"DOCX提取失败: {file_path} - {e}")
        return None


def extract_pptx(file_path):
    try:
        pres = pptx.Presentation(file_path)
        text = ""
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"PPTX提取失败: {file_path} - {e}")
        return None


def extract_pdf(file_path):
    try:
        doc = fitz.open(file_path)
        text = '\n'.join(page.get_text() for page in doc)
        if not text.strip():
            raise ValueError("PDF无文本内容，尝试OCR")
        return text
    except Exception as e:
        print(f"PDF原生提取失败，尝试OCR：{file_path} - {e}")
        return ocr_pdf(file_path)


def extract_wbd(file_path):
    # .wbd 是希沃白板文件，通常是压缩包或 XML 格式，需要特殊解析器
    print(f"WBD文件暂不支持内容提取：{file_path}")
    return ""


def ocr_image(image_path):
    try:
        img = cv2.imread(image_path)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        _, binary = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)
        text = pytesseract.image_to_string(binary, lang='chi_sim')
        return text
    except Exception as e:
        print(f"OCR识别失败: {image_path} - {e}")
        return ""


def ocr_pdf(pdf_path):
    """将 PDF 转为图像后 OCR 提取文本"""
    try:
        doc = fitz.open(pdf_path)
        temp_dir = tempfile.mkdtemp()
        texts = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=200)
            img_path = os.path.join(temp_dir, f"page_{page_num}.png")
            pix.save(img_path)
            texts.append(ocr_image(img_path))

        shutil.rmtree(temp_dir)
        return '\n'.join(texts)
    except Exception as e:
        print(f"OCR转换失败: {pdf_path} - {e}")
        return ""


# ==================== 加载模型 ====================
try:
    clf = joblib.load("subject_classifier.pkl")
    vectorizer = joblib.load("tfidf_vectorizer.pkl")
except Exception as e:
    print("⚠️ 模型文件未找到，请先运行训练脚本！", e)
    exit(1)


# ==================== 多维分类逻辑 ====================
def classify_file(file_path):
    filename = os.path.basename(file_path)
    usb_label = get_usb_label(file_path)
    content = extract_content(file_path)

    scores = defaultdict(int)

    # 文件名匹配
    for subject, keywords in SUBJECT_KEYWORDS.items():
        for keyword in keywords:
            if re.search(keyword, filename, re.IGNORECASE):
                scores[subject] += 1

    # U盘标签匹配（优先级最高）
    if usb_label in USB_LABEL_MAP:
        scores[USB_LABEL_MAP[usb_label]] += 3

    # 内容匹配
    if content and len(content) > 30:
        X = vectorizer.transform([content])
        ai_subject = clf.predict(X)[0]
        scores[ai_subject] += 2

    # 决策逻辑
    if scores:
        max_score = max(scores.values())
        candidates = [k for k, v in scores.items() if v == max_score]
        return candidates[0] if len(candidates) == 1 else None
    return None


# ==================== 提取文件内容 ====================
def extract_content(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    elif ext == ".pdf":
        return extract_pdf(file_path)
    elif ext == ".wbd":
        return extract_wbd(file_path)
    else:
        return ""


# ==================== 文件处理缓存 ====================
processed_files = set()
processed_lock = threading.Lock()


# ==================== 文件监控 ====================
class FileHandler(FileSystemEventHandler):
    def on_created(self, event):
        if event.is_directory:
            return

        file_path = event.src_path
        _, ext = os.path.splitext(file_path)

        if ext.lower() not in SUPPORTED_EXTS:
            return

        with processed_lock:
            if file_path in processed_files:
                print(f"[{time.strftime('%H:%M:%S')}] 跳过已处理文件：{file_path}")
                return
            processed_files.add(file_path)

        threading.Thread(target=self.delayed_classification, args=(file_path,), daemon=True).start()

    def delayed_classification(self, file_path):
        print(f"[{time.strftime('%H:%M:%S')}] 等待 {DELAY_SECONDS} 秒后尝试识别文件：{file_path}")
        time.sleep(DELAY_SECONDS)

        if not os.path.exists(file_path):
            print(f"[{time.strftime('%H:%M:%S')}] 文件不存在：{file_path}")
            return

        try:
            with open(file_path, 'rb'):
                pass
        except Exception as e:
            print(f"[{time.strftime('%H:%M:%S')}] 文件无法访问（仍在复制中？）：{file_path} - {e}")
            return

        subject = classify_file(file_path)
        if subject:
            dest_folder = os.path.join(OUTPUT_BASE_FOLDER, subject)
            ensure_folder_exists(dest_folder)
            dest_path = os.path.join(dest_folder, os.path.basename(file_path))
            try:
                shutil.move(file_path, dest_path)
                print(f"[{time.strftime('%H:%M:%S')}] 文件 {file_path} 分类为：{subject}，已移动至 {dest_folder}")
            except Exception as e:
                print(f"[{time.strftime('%H:%M:%S')}] 文件移动失败：{file_path} -> {dest_folder} - {e}")
        else:
            notification.notify(
                title="文件分类提醒",
                message=f"无法识别文件类型或内容太少：{file_path}",
                timeout=5
            )


# ==================== 主程序入口 ====================
if __name__ == "__main__":
    print("=== 本地AI文件分类器 & 自动归类工具 ===")
    print(f"正在监控目录：{WATCH_FOLDER}")
    print(f"归类目标目录：{OUTPUT_BASE_FOLDER}")
    print("按 Ctrl+C 停止监控")

    observer = Observer()
    observer.schedule(FileHandler(), path=WATCH_FOLDER, recursive=False)
    observer.start()

    try:
        while True:
            time.sleep(1)
    except KeyboardInterrupt:
        print("\n停止监控...")
        observer.stop()
    observer.join()