import os
import re
from docx import Document
from pptx import Presentation
import PyPDF2
from PIL import Image
import pytesseract
import cv2
import tempfile
import fitz  # PyMuPDF
import shutil
import json

# 设置 Tesseract 路径
CONFIG_FILE = "config.json"
try:
    with open(CONFIG_FILE, "r", encoding="utf-8") as f:
        config = json.load(f)
except FileNotFoundError:
    print("⚠️ 未找到配置文件，请先运行配置工具")
    exit(1)
pytesseract.pytesseract.tesseract_cmd = config["TESSERACT_PATH"]

def clean_text(text):
    return re.sub(r'\s+', ' ', text).strip()


def extract_docx(file_path):
    try:
        doc = Document(file_path)
        return '\n'.join(paragraph.text for paragraph in doc.paragraphs)
    except Exception as e:
        print(f"DOCX提取失败: {file_path} - {e}")
        return ""


def extract_pptx(file_path):
    try:
        pres = Presentation(file_path)
        text = ""
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"PPTX提取失败: {file_path} - {e}")
        return ""


def extract_pdf(file_path):
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            return '\n'.join(page.extract_text() for page in reader.pages)
    except Exception as e:
        print(f"PDF提取失败（尝试OCR）: {file_path} - {e}")
        return ocr_pdf(file_path)


def ocr_pdf(pdf_path):
    """将 PDF 转为图像后 OCR 提取文本"""
    doc = fitz.open(pdf_path)
    temp_dir = tempfile.mkdtemp()
    texts = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=200)
        img_path = os.path.join(temp_dir, f"page_{page_num}.png")
        pix.save(img_path)
        texts.append(ocr_image(img_path))

    shutil.rmtree(temp_dir)
    return '\n'.join(texts)


def ocr_image(image_path):
    """对图像进行 OCR 识别"""
    try:
        img = cv2.imread(image_path)
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        _, binary = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)
        # OCR 识别
        text = pytesseract.image_to_string(binary, lang='chi_sim')
        return text
    except Exception as e:
        print(f"OCR识别失败: {image_path} - {e}")
        return ""


def extract_content(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    elif ext == ".pdf":
        return extract_pdf(file_path)
    else:
        return ""


def process_folder(data_dir, output_dir):
    categories = [d for d in os.listdir(data_dir) if os.path.isdir(os.path.join(data_dir, d))]

    for category in categories:
        input_folder = os.path.join(data_dir, category)
        output_folder = os.path.join(output_dir, category)
        os.makedirs(output_folder, exist_ok=True)

        print(f"正在处理科目：{category}")
        count = 0
        for filename in os.listdir(input_folder):
            file_path = os.path.join(input_folder, filename)
            content = extract_content(file_path)
            cleaned = clean_text(content)
            if cleaned and len(cleaned) >= 30:
                out_file = os.path.join(output_folder, f"{os.path.splitext(filename)[0]}.txt")
                with open(out_file, "w", encoding="utf-8", errors="ignore") as f:
                    f.write(cleaned)
                count += 1
            else:
                print(f"跳过无效文件或内容太少：{file_path}")
        print(f"已提取 {count} 个文件到 {output_folder}")


if __name__ == "__main__":
    data_dir = input("请输入训练数据根目录（如 train_data）: ").strip()
    output_dir = input("请输入输出目录（如 extracted_texts）: ").strip()
    process_folder(data_dir, output_dir)